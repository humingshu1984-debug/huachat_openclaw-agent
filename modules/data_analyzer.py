import pandas as pd
import yaml
import os
from pptx import Presentation

class DataAnalyzer:
    def __init__(self, config_path="config/config.yaml"):
        with open(config_path, 'r', encoding='utf-8') as f:
            self.config = yaml.safe_load(f)
        self.template_groups = []
        self._load_template_groups()

    def _normalize_col_name(self, col):
        return "".join(str(col).replace("\u3000", " ").split()).strip()

    def detect_status_col(self, df):
        candidates = []
        for col in df.columns:
            norm = self._normalize_col_name(col)
            if not norm:
                continue
            if norm in {"完成状态", "学习状态"}:
                return col
            if "完成" in norm and "状态" in norm:
                candidates.append(col)
            elif "学习" in norm and "状态" in norm:
                candidates.append(col)
        return candidates[0] if candidates else df.columns[-1]

    def detect_account_status_col(self, df):
        candidates = []
        for col in df.columns:
            norm = self._normalize_col_name(col)
            if not norm:
                continue
            if norm in {"账号状态", "账户状态"}:
                return col
            if ("账号" in norm or "账户" in norm) and "状态" in norm:
                candidates.append(col)
        return candidates[0] if candidates else None

    def _filter_disabled_accounts(self, df):
        col = self.detect_account_status_col(df)
        if not col:
            return df
        status = df[col].astype(str).str.strip()
        return df[status != "已禁用"].copy()

    def _load_template_groups(self):
        """从 PPT 模板中加载小组名单"""
        template_path = "data/PPT模板.pptx"
        if os.path.exists(template_path):
            try:
                prs = Presentation(template_path)
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if shape.has_table:
                            table = shape.table
                            groups = []
                            for i, row in enumerate(table.rows):
                                if i == 0: continue # 跳过表头
                                group_name = row.cells[0].text_frame.text.strip().replace(' ', '').replace('\n', '')
                                if group_name and group_name != '总计':
                                    groups.append(group_name)
                            self.template_groups = groups
                            print(f"从 PPT 模板识别到 {len(groups)} 个小组: {groups}")
                            return
            except Exception as e:
                print(f"读取 PPT 模板小组名单失败: {e}")
        
        self.template_groups = ['非洲区', '中东北非区', '欧亚区', '大洋洲区', '英国区', '南洋区', '东北欧区', '拉美区', '东南亚区', '巴西区', '墨西哥区', '西欧区', '国际合作部']

    def parse_department_level(self, dept_str, level):
        if pd.isna(dept_str) or not isinstance(dept_str, str):
            return "未知"
        parts = [p.strip() for p in dept_str.split('->') if str(p).strip() != ""]
        if len(parts) >= level:
            return parts[level - 1].strip()
        if parts:
            return parts[-1].strip()
        return "未知"

    def parse_department(self, dept_str):
        """解析部门字符串，提取第三级部门"""
        return self.parse_department_level(dept_str, 3)

    def analyze_learning_progress(self, raw_data_path, group_filter=None):
        """分析单份学习进度数据"""
        df = pd.read_excel(raw_data_path)
        df = self._filter_disabled_accounts(df)
        
        # 1. 提取三级部门
        dept_col = '部门' if '部门' in df.columns else '所属小组'
        df['三级部门'] = df[dept_col].apply(lambda x: self.parse_department_level(x, 3))
        df['四级部门'] = df[dept_col].apply(lambda x: self.parse_department_level(x, 4))
        df['分析部门'] = df['三级部门']
        
        if group_filter:
            df = df[df['分析部门'].isin(group_filter)].copy()

        # 2. 统计已完成情况
        status_col = self.detect_status_col(df)

        # 计算每个部门的学员总数和已完成人数
        grouped = df.groupby('分析部门').agg({
            status_col: [
                ('总人数', 'count'),
                ('已完成人数', lambda x: (x == '已完成').sum())
            ]
        })
        grouped.columns = grouped.columns.get_level_values(1)
        grouped = grouped.reset_index()
        grouped = grouped.rename(columns={'分析部门': '小组'})

        return grouped, df

    def build_group_summary(self, df, group_col, status_col):
        grouped = df.groupby(group_col).agg({
            status_col: [
                ('总人数', 'count'),
                ('已完成人数', lambda x: (x == '已完成').sum())
            ]
        })
        grouped.columns = grouped.columns.get_level_values(1)
        grouped = grouped.reset_index()
        grouped = grouped.rename(columns={group_col: '小组'})
        return grouped

    def sanitize_sheet_name(self, name):
        safe = str(name)
        for ch in ['\\', '/', '?', '*', '[', ']', ':']:
            safe = safe.replace(ch, ' ')
        safe = safe.strip()
        if not safe:
            safe = "Sheet"
        if len(safe) > 31:
            safe = safe[:31]
        return safe

    def save_analysis_excel(self, main_df, output_path, extra_sheets=None):
        os.makedirs(os.path.dirname(output_path), exist_ok=True)
        with pd.ExcelWriter(output_path, engine='openpyxl') as writer:
            main_df.to_excel(writer, index=False, sheet_name='汇总')
            if extra_sheets:
                for sheet_name, df in extra_sheets.items():
                    safe_name = self.sanitize_sheet_name(sheet_name)
                    df.to_excel(writer, index=False, sheet_name=safe_name)
        return output_path

    def save_uncompleted_lists(self, all_raw_dfs, output_dir, group_filter=None, exclude_filter=None):
        """为每个小组生成未完成学员清单"""
        # 合并所有原始数据
        combined_raw = pd.concat(all_raw_dfs, ignore_index=True)
        
        # 确定状态列
        status_col = self.detect_status_col(combined_raw)

        # 过滤未完成学员
        uncompleted_df = combined_raw[combined_raw[status_col] != '已完成'].copy()
        if group_filter and '分析部门' in uncompleted_df.columns:
            uncompleted_df = uncompleted_df[uncompleted_df['分析部门'].isin(group_filter)].copy()
        if exclude_filter and '分析部门' in uncompleted_df.columns:
            uncompleted_df = uncompleted_df[~uncompleted_df['分析部门'].isin(exclude_filter)].copy()
        
        # 按小组分组并保存
        uncompleted_dir = os.path.join(output_dir, "uncompleted_lists")
        os.makedirs(uncompleted_dir, exist_ok=True)
        
        saved_files = []
        for group_name, group_data in uncompleted_df.groupby('分析部门'):
            if not group_name or group_name == "未知":
                continue
            if group_filter and group_name not in group_filter:
                continue
            if exclude_filter and group_name in exclude_filter:
                continue
            
            # 移除中间辅助列
            output_cols = [c for c in group_data.columns if c != '分析部门']
            file_name = f"{group_name}未完成学员清单.xlsx"
            file_path = os.path.join(uncompleted_dir, file_name)
            
            group_data[output_cols].to_excel(file_path, index=False)
            saved_files.append(file_path)
            
        return saved_files

    def merge_results(self, results_list, previous_df=None):
        """汇总多份数据结果并进行排名"""
        # 合并 data1 和 data2 的原始统计数据
        combined = pd.concat(results_list, ignore_index=True)
        
        # 按小组汇总（累加人数）
        final_grouped = combined.groupby('小组').agg({
            '总人数': 'sum',
            '已完成人数': 'sum'
        }).reset_index()

        # 计算完成率
        final_grouped['学习完成率'] = (final_grouped['已完成人数'] / final_grouped['总人数'] * 100).round(2)

        # 3. 按照完成率排名
        final_df = final_grouped.sort_values(by='学习完成率', ascending=False).reset_index(drop=True)

        # 添加排名列
        final_df['本轮排名'] = final_df.index + 1
        
        # 4. 对比上一轮数据
        final_df['上一轮排名'] = ""
        final_df['上一轮学习完成率'] = ""

        if previous_df is not None:
            # 建立小组到排名/完成率的映射（排除总计行）
            prev_data = previous_df[previous_df['小组'] != '总计'].copy()
            # 统一小组名称格式（去除空格）
            prev_data['小组_key'] = prev_data['小组'].astype(str).str.replace(' ', '').str.replace('\n', '')
            
            # 创建映射字典
            rank_map = dict(zip(prev_data['小组_key'], prev_data['本轮排名']))
            rate_map = dict(zip(prev_data['小组_key'], prev_data['学习完成率']))

            # 匹配到当前数据中
            def get_prev_rank(group_name):
                key = str(group_name).replace(' ', '').replace('\n', '')
                return rank_map.get(key, "")

            def get_prev_rate(group_name):
                key = str(group_name).replace(' ', '').replace('\n', '')
                return rate_map.get(key, "")

            final_df['上一轮排名'] = final_df['小组'].apply(get_prev_rank)
            final_df['上一轮学习完成率'] = final_df['小组'].apply(get_prev_rate)

        # 5. 计算总计行
        total_sum = final_df[['总人数', '已完成人数']].sum()
        total_rate = (total_sum['已完成人数'] / total_sum['总人数'] * 100).round(2)
        
        # 处理总计行的上一轮数据
        prev_total_rank = "-"
        prev_total_rate = "-"
        if previous_df is not None:
            prev_total_row = previous_df[previous_df['小组'] == '总计']
            if not prev_total_row.empty:
                # 在历史数据中，当时的总计完成率存放在 '学习完成率' 列
                prev_total_rate = prev_total_row.iloc[0]['学习完成率']

        total_row = pd.DataFrame({
            '小组': ['总计'],
            '总人数': [total_sum['总人数']],
            '已完成人数': [total_sum['已完成人数']],
            '学习完成率': [total_rate],
            '本轮排名': ['-'],
            '上一轮排名': [prev_total_rank],
            '上一轮学习完成率': [prev_total_rate]
        })

        final_df = pd.concat([final_df, total_row], ignore_index=True)

        # 格式化输出百分比
        output_df = final_df.copy()
        output_df['学习完成率'] = output_df['学习完成率'].apply(lambda x: f"{x}%" if x != '-' else x)
        
        # 调整列顺序以匹配汇报要求
        cols = ['小组', '总人数', '已完成人数', '学习完成率', '本轮排名', '上一轮排名', '上一轮学习完成率']
        output_df = output_df[cols]

        return output_df

    def save_analysis(self, result_df, output_path):
        """保存分析结果"""
        os.makedirs(os.path.dirname(output_path), exist_ok=True)
        result_df.to_excel(output_path, index=False)
        return output_path
