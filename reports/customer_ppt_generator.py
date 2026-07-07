import os
import time
import yaml
import pandas as pd
from copy import deepcopy
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


class CustomerPPTGenerator:
    def __init__(self, config_path: str):
        with open(config_path, "r", encoding="utf-8") as f:
            self.config = yaml.safe_load(f) or {}
        self.paths = self.config.get("paths", {}) or {}
        self.ppt_cfg = self.config.get("ppt", {}) or {}

    def _slide_text(self, slide):
        parts = []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                t = (shape.text_frame.text or "").strip()
                if t:
                    parts.append(t)
        return "\n".join(parts)

    def _contains_any(self, text, keywords):
        if not text:
            return False
        for k in keywords or []:
            if k and str(k) in text:
                return True
        return False

    def _ensure_table(self, slide, table_shape, rows, cols, force_rebuild=False):
        table = table_shape.table
        if not force_rebuild and len(table.rows) >= rows and len(table.columns) >= cols:
            return table_shape, table
        left, top, width, height = table_shape.left, table_shape.top, table_shape.width, table_shape.height
        sp = table_shape._element
        sp.getparent().remove(sp)
        new_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
        return new_shape, new_shape.table

    def _apply_cell_style(self, cell, is_header=False, preserve_template_style=True, font_size_pt=7, remove_fill=False):
        if remove_fill:
            try:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(255, 255, 255)
            except Exception:
                pass
        for paragraph in cell.text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            for run in paragraph.runs:
                run.font.size = Pt(font_size_pt)
                run.font.name = "微软雅黑"
                if not preserve_template_style:
                    run.font.color.rgb = RGBColor(0, 0, 0)
                    if is_header:
                        run.font.bold = True
                elif is_header:
                    run.font.bold = True

    def _clone_slide(self, prs, source_slide):
        new_slide = prs.slides.add_slide(prs.slide_layouts[6])
        for shape in list(new_slide.shapes):
            sp = shape._element
            sp.getparent().remove(sp)
        for shape in source_slide.shapes:
            new_el = deepcopy(shape._element)
            new_slide.shapes._spTree.insert_element_before(new_el, "p:extLst")
        return new_slide

    def _find_first_table_shape(self, slide):
        for shape in slide.shapes:
            if getattr(shape, "has_table", False):
                return shape
        return None

    def _find_summary_box(self, slide):
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            if int(shape.left) >= int(Inches(8.0)) and int(shape.top) <= int(Inches(2.0)):
                return shape
        return None

    def _replace_text_with_style(self, shape, text, font_size_pt, bold=False, color_rgb=(0, 0, 0)):
        if not getattr(shape, "has_text_frame", False):
            return
        tf = shape.text_frame
        tf.clear()
        p = tf.paragraphs[0] if tf.paragraphs else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = text
        run.font.size = Pt(font_size_pt)
        run.font.bold = bold
        run.font.name = "微软雅黑"
        run.font.color.rgb = RGBColor(*color_rgb)

    def _set_slide_titles(self, slide, small_title, main_title):
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            text = (shape.text_frame.text or "").strip()
            if "方案-极客行动" in text:
                self._replace_text_with_style(shape, small_title, 11, False, (128, 128, 128))
            elif ("极客“掘金”" in text and "客户方案技术营销" in text and "培训覆盖" not in text) or ("客户培训覆盖" in text and "整体情况" in text):
                self._replace_text_with_style(shape, main_title, 14, True, (0, 0, 0))

    def _fill_table(self, slide, table_shape, df: pd.DataFrame, preserve_template_style=True, font_size_pt=7, remove_fill=False, force_rebuild=False):
        if df is None or df.empty:
            return
        df2 = df.copy()
        df2 = df2.fillna("")
        first_col = str(df2.columns[0])
        df2 = df2[df2[first_col].astype(str).str.strip() != ""].copy()
        rows = len(df2) + 1
        cols = len(df2.columns)
        table_shape, table = self._ensure_table(slide, table_shape, rows, cols, force_rebuild=force_rebuild)

        headers = list(df2.columns)
        for c, h in enumerate(headers):
            cell = table.cell(0, c)
            cell.text = str(h)
            self._apply_cell_style(cell, is_header=True, preserve_template_style=preserve_template_style, font_size_pt=font_size_pt, remove_fill=remove_fill)

        for r_idx, (_, row) in enumerate(df2.iterrows(), start=1):
            for c_idx, h in enumerate(headers):
                cell = table.cell(r_idx, c_idx)
                cell.text = str(row[h])
                self._apply_cell_style(cell, is_header=False, preserve_template_style=preserve_template_style, font_size_pt=font_size_pt, remove_fill=remove_fill)

        for r_idx in range(len(df2) + 1, len(table.rows)):
            for c_idx in range(cols):
                cell = table.cell(r_idx, c_idx)
                cell.text = ""
                self._apply_cell_style(cell, is_header=False, preserve_template_style=preserve_template_style, font_size_pt=font_size_pt, remove_fill=remove_fill)

    def _clear_textboxes_containing(self, slide, keywords):
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            try:
                text = shape.text_frame.text or ""
            except Exception:
                continue
            if any(k in text for k in keywords or []):
                try:
                    shape.text_frame.clear()
                except Exception:
                    try:
                        shape.text_frame.text = ""
                    except Exception:
                        pass

    def _format_pct(self, v):
        try:
            return f"{float(v):.2f}%"
        except Exception:
            return str(v)

    def _group_by_tier(self, df, name_col, coverage_col=None):
        leaders = []
        mid = []
        low = []
        if coverage_col is None:
            coverage_col = "coverage_rate" if "coverage_rate" in df.columns else "covered_ratio"
        for _, row in df.iterrows():
            item = f"{row[name_col]}（{self._format_pct(row[coverage_col])}）"
            rate = float(row[coverage_col])
            if rate >= 15:
                leaders.append(item)
            elif rate >= 5:
                mid.append(item)
            else:
                low.append(item)
        return leaders, mid, low

    def _add_summary_text(self, slide, summary_df: pd.DataFrame, name_col="region_l2", scope_name=None, item_label="大区", coverage_col=None, cert_col=None):
        if summary_df is None or summary_df.empty:
            return
        df = summary_df.copy()
        if coverage_col is None:
            coverage_col = "coverage_rate" if "coverage_rate" in df.columns else "covered_ratio"
        if cert_col is None:
            cert_col = "certified_rate" if "certified_rate" in df.columns else "certified_ratio"
        df = df.sort_values(by=coverage_col, ascending=False).reset_index(drop=True)
        top = df.iloc[0]
        bottom = df.iloc[-1]
        leaders, mid, low = self._group_by_tier(df, name_col, coverage_col=coverage_col)

        self._clear_textboxes_containing(slide, ["极客“掘金”-客户技术营销培训覆盖", "本月极客掘金覆盖客户数", "客户覆盖：", "10月客户极客掘金圈定覆盖新增"])
        summary_box = self._find_summary_box(slide)
        if summary_box is None:
            summary_box = slide.shapes.add_textbox(Inches(8.2), Inches(0.95), Inches(4.0), Inches(2.2))
        tf = summary_box.text_frame
        tf.word_wrap = True
        try:
            tf.clear()
        except Exception:
            pass

        p0 = tf.paragraphs[0] if tf.paragraphs else tf.add_paragraph()
        p0.clear()
        title_run = p0.add_run()
        title_run.text = f"{scope_name}关键结论" if scope_name else "整体关键结论"
        title_run.font.bold = True
        title_run.font.size = Pt(15)
        title_run.font.name = "微软雅黑"
        title_run.font.color.rgb = RGBColor(192, 0, 0)

        lines = []
        if len(df) == 1:
            target = (
                f"{scope_name}下当前仅有 1 个{item_label}纳入统计，为 {top[name_col]}，已覆盖比例为 {self._format_pct(top[coverage_col])}。"
                if scope_name
                else f"当前仅有 1 个{item_label}纳入统计，为 {top[name_col]}，已覆盖比例为 {self._format_pct(top[coverage_col])}。"
            )
            lines.append(f"1. {target}")
        else:
            prefix = f"{scope_name}下各{item_label}" if scope_name else f"各{item_label}"
            lines.append(
                f"1. {prefix}按已覆盖比例排序后，最高为 {top[name_col]}（{self._format_pct(top[coverage_col])}），最低为 {bottom[name_col]}（{self._format_pct(bottom[coverage_col])}）。"
            )

        if leaders:
            summary_parts = [f"进展较好的{item_label}：{'、'.join(leaders[:3])}"]
            if low:
                summary_parts.append(f"需重点跟进的{item_label}：{'、'.join(low[:3])}")
            elif mid:
                summary_parts.append(f"仍需提升的{item_label}：{'、'.join(mid[:3])}")
            lines.append(f"2. {'；'.join(summary_parts)}。")
        elif low:
            lines.append(f"2. 当前需重点跟进的{item_label}：{'、'.join(low[:3])}。")

        cert_top = df.sort_values(by=cert_col, ascending=False).iloc[0]
        lines.append(f"3. 认证推进最快的{item_label}是 {cert_top[name_col]}（{self._format_pct(cert_top[cert_col])}）。")
        for t in lines:
            p = tf.add_paragraph()
            p.clear()
            run = p.add_run()
            run.text = t
            run.font.size = Pt(11.5)
            run.font.name = "微软雅黑"
            run.font.color.rgb = RGBColor(0, 0, 0)

    def generate_report(self, dfs: dict, template_path: str, output_path: str):
        if not template_path:
            template_path = self.paths.get("template_ppt")
        if not os.path.exists(template_path):
            raise FileNotFoundError(f"找不到 PPT 模板：{template_path}")

        prs = Presentation(template_path)

        coverage_keywords = (self.ppt_cfg.get("slide_title_candidates", {}) or {}).get("coverage", []) or []
        certification_keywords = (self.ppt_cfg.get("slide_title_candidates", {}) or {}).get("certification", []) or []

        tables = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_table:
                    tables.append((slide, shape))

        if not tables:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            shape = slide.shapes.add_table(2, 2, Inches(0.6), Inches(1.4), Inches(8.0), Inches(4.5))
            tables.append((slide, shape))

        region_summary = dfs.get("region_summary_df")
        ppt_region_summary = dfs.get("ppt_region_summary_df")
        subregion_summary = dfs.get("subregion_summary_df")
        ppt_subregion_summary_map = dfs.get("ppt_subregion_summary_map") or {}
        region_l3_summary = dfs.get("region_level3_summary_df")
        cert_type_count = dfs.get("cert_type_count_df")
        coverage_x_cert = dfs.get("coverage_x_cert_df")
        top_uncovered = dfs.get("top_uncovered_customers_df")

        assigned = []
        for slide, shape in tables:
            text = self._slide_text(slide)
            if self._contains_any(text, coverage_keywords):
                assigned.append(("region_summary", slide, shape))
            elif self._contains_any(text, certification_keywords):
                assigned.append(("cert_type_count", slide, shape))
            else:
                assigned.append((None, slide, shape))

        coverage_target = None
        for name, slide, shape in assigned:
            if name == "region_summary":
                coverage_target = (slide, shape)
                break
        if coverage_target is None and assigned:
            coverage_target = (assigned[0][1], assigned[0][2])

        region_slide_specs = []
        if coverage_target is not None and region_summary is not None and not region_summary.empty:
            for region_name in region_summary["region_l2"].tolist():
                if region_name not in ppt_subregion_summary_map:
                    continue
                region_slide_specs.append((region_name, self._clone_slide(prs, coverage_target[0])))

        used = set()
        for name, slide, shape in assigned:
            if name == "region_summary" and "region_summary" not in used:
                self._set_slide_titles(slide, "方案-极客行动 | 整体情况", "客户培训覆盖整体情况")
                self._fill_table(
                    slide,
                    shape,
                    ppt_region_summary if ppt_region_summary is not None else region_summary,
                    preserve_template_style=True,
                    font_size_pt=7,
                    remove_fill=False,
                    force_rebuild=False,
                )
                self._add_summary_text(slide, region_summary, name_col="region_l2", item_label="大区")
                used.add("region_summary")
            elif name == "cert_type_count" and "cert_type_count" not in used:
                self._fill_table(slide, shape, cert_type_count, preserve_template_style=True, font_size_pt=7, remove_fill=False, force_rebuild=False)
                used.add("cert_type_count")

        for name, slide, shape in assigned:
            if name is None:
                if "region_summary" not in used:
                    self._set_slide_titles(slide, "方案-极客行动 | 整体情况", "客户培训覆盖整体情况")
                    self._fill_table(
                        slide,
                        shape,
                        ppt_region_summary if ppt_region_summary is not None else region_summary,
                        preserve_template_style=True,
                        font_size_pt=7,
                        remove_fill=False,
                        force_rebuild=False,
                    )
                    self._add_summary_text(slide, region_summary, name_col="region_l2", item_label="大区")
                    used.add("region_summary")
                elif "cert_type_count" not in used:
                    self._fill_table(slide, shape, cert_type_count, preserve_template_style=True, font_size_pt=7, remove_fill=False, force_rebuild=False)
                    used.add("cert_type_count")
                elif "coverage_x_cert" not in used:
                    self._fill_table(slide, shape, coverage_x_cert, preserve_template_style=True, font_size_pt=7, remove_fill=False, force_rebuild=False)
                    used.add("coverage_x_cert")
                elif "region_l3_summary" not in used:
                    self._fill_table(slide, shape, region_l3_summary, preserve_template_style=True, font_size_pt=7, remove_fill=False, force_rebuild=False)
                    used.add("region_l3_summary")
                elif "top_uncovered" not in used:
                    self._fill_table(slide, shape, top_uncovered, preserve_template_style=True, font_size_pt=7, remove_fill=False, force_rebuild=False)
                    used.add("top_uncovered")

        for region_name, slide in region_slide_specs:
            table_shape = self._find_first_table_shape(slide)
            if table_shape is None:
                continue
            self._set_slide_titles(slide, f"方案-极客行动 | {region_name}", f"{region_name}代表处客户培训覆盖情况")
            self._fill_table(
                slide,
                table_shape,
                ppt_subregion_summary_map.get(region_name),
                preserve_template_style=False,
                font_size_pt=7,
                remove_fill=True,
                force_rebuild=True,
            )
            region_detail_df = None
            if subregion_summary is not None and not subregion_summary.empty:
                region_detail_df = subregion_summary[subregion_summary["region_l2"] == region_name].copy()
            self._add_summary_text(
                slide,
                region_detail_df,
                name_col="region_l3",
                scope_name=region_name,
                item_label="代表处",
            )

        software_region_summary = dfs.get("software_region_summary_df")
        if software_region_summary is not None and isinstance(software_region_summary, pd.DataFrame) and not software_region_summary.empty:
            base_slide = coverage_target[0] if coverage_target is not None else (assigned[0][1] if assigned else None)
            if base_slide is None:
                software_slide = prs.slides.add_slide(prs.slide_layouts[1])
                table_shape = software_slide.shapes.add_table(2, 2, Inches(0.6), Inches(1.4), Inches(8.0), Inches(4.5))
            else:
                software_slide = self._clone_slide(prs, base_slide)
                table_shape = self._find_first_table_shape(software_slide)

            if table_shape is not None:
                df_ppt = software_region_summary.copy()
                if "已覆盖比例" in df_ppt.columns:
                    df_ppt["已覆盖比例"] = df_ppt["已覆盖比例"].map(lambda x: "" if pd.isna(x) else f"{float(x):.2f}%")
                if "Software认证通过比例" in df_ppt.columns:
                    df_ppt["Software认证通过比例"] = df_ppt["Software认证通过比例"].map(lambda x: "" if pd.isna(x) else f"{float(x):.2f}%")
                if "DHIA认证通过比例" in df_ppt.columns:
                    df_ppt["DHIA认证通过比例"] = df_ppt["DHIA认证通过比例"].map(lambda x: "" if pd.isna(x) else f"{float(x):.2f}%")

                keep_cols = [c for c in ["大区", "NP客户数量", "已覆盖比例", "Software认证通过个数", "Software认证通过比例", "Software认证通过比例排名", "DHIA认证通过比例"] if c in df_ppt.columns]
                if keep_cols:
                    df_ppt = df_ppt[keep_cols]

                self._set_slide_titles(software_slide, "方案-极客行动 | Software", "DHIA - Software 认证情况")
                self._fill_table(
                    software_slide,
                    table_shape,
                    df_ppt,
                    preserve_template_style=True,
                    font_size_pt=7,
                    remove_fill=False,
                    force_rebuild=True,
                )
                if region_summary is not None and not region_summary.empty and "software_certified_ratio" in region_summary.columns:
                    self._add_summary_text(
                        software_slide,
                        region_summary,
                        name_col="region_l2",
                        item_label="大区",
                        coverage_col="covered_ratio",
                        cert_col="software_certified_ratio",
                    )
                top_uncert = dfs.get("software_top_uncertified_customers_df")
                if top_uncert is not None and isinstance(top_uncert, pd.DataFrame) and not top_uncert.empty:
                    top_uncert = top_uncert.copy()
                    names = top_uncert["客户名称"].astype(str).str.strip().replace("nan", "").tolist() if "客户名称" in top_uncert.columns else []
                    names = [n for n in names if n]
                    if names:
                        box = software_slide.shapes.add_textbox(Inches(0.6), Inches(6.2), Inches(12.0), Inches(0.9))
                        tf = box.text_frame
                        tf.word_wrap = True
                        try:
                            tf.clear()
                        except Exception:
                            pass
                        p0 = tf.paragraphs[0] if tf.paragraphs else tf.add_paragraph()
                        p0.clear()
                        r0 = p0.add_run()
                        r0.text = f"Top未认证客户（已覆盖但未通过Software认证）：{'、'.join(names[:10])}"
                        r0.font.size = Pt(11)
                        r0.font.bold = True
                        r0.font.name = "微软雅黑"
                        r0.font.color.rgb = RGBColor(0, 0, 0)

        os.makedirs(os.path.dirname(output_path), exist_ok=True)
        try:
            prs.save(output_path)
        except PermissionError:
            base, ext = os.path.splitext(output_path)
            output_path = f"{base}_{time.strftime('%Y%m%d_%H%M%S')}{ext}"
            prs.save(output_path)
        return output_path
