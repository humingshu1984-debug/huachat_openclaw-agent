from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import pandas as pd
import yaml
import os


class PPTGenerator:
    def __init__(self, config_path="config/config.yaml"):
        with open(config_path, "r", encoding="utf-8") as f:
            self.config = yaml.safe_load(f)

    def generate_report(self, result_df, template_path=None, output_path="outputs/final_weekly_report.pptx"):
        if template_path is None:
            template_path = self.config["paths"]["template_ppt"]

        if os.path.exists(template_path):
            prs = Presentation(template_path)
        else:
            print(f"⚠️ 模板文件 {template_path} 不存在，将创建一个空的 PPT。")
            prs = Presentation()

        analysis_data = result_df[result_df["小组"] != "总计"].copy()
        total_rows = result_df[result_df["小组"] == "总计"]
        if total_rows.empty:
            total_sum = analysis_data[["总人数", "已完成人数"]].sum()
            total_rate = (total_sum["已完成人数"] / total_sum["总人数"] * 100).round(2) if total_sum["总人数"] else 0
            total_info = {"总人数": int(total_sum["总人数"]), "学习完成率": f"{total_rate}%"}
        else:
            total_info = total_rows.iloc[0]

        slide = None
        for s in prs.slides:
            if any(shape.has_table for shape in s.shapes):
                slide = s
                break
        if not slide:
            slide = prs.slides.add_slide(prs.slide_layouts[1])

        table = None
        table_shape = None
        for shape in slide.shapes:
            if shape.has_table:
                table_shape = shape
                table = shape.table
                if len(table.rows) < len(result_df) + 1 or len(table.columns) < len(result_df.columns):
                    print(
                        f"⚠️ 模板表格大小不足 ({len(table.rows)}x{len(table.columns)})，正在重新创建以适配数据 ({len(result_df)+1}x{len(result_df.columns)})"
                    )
                    left, top, width, height = table_shape.left, table_shape.top, table_shape.width, table_shape.height
                    sp = table_shape._element
                    sp.getparent().remove(sp)
                    rows, cols = len(result_df) + 1, len(result_df.columns)
                    shape = slide.shapes.add_table(rows, cols, left, top, width, height)
                    table = shape.table
                break

        if not table:
            rows, cols = len(result_df) + 1, len(result_df.columns)
            left, top, width, height = Inches(0.5), Inches(1.2), Inches(6.5), Inches(0.5)
            shape = slide.shapes.add_table(rows, cols, left, top, width, height)
            table = shape.table

        self._clear_textboxes_containing(slide, ["能力地图-内训赋能进展"])

        headers = result_df.columns.tolist()
        num_rows = len(table.rows)
        num_cols = len(table.columns)
        for i, h in enumerate(headers):
            if i >= num_cols:
                continue
            cell = table.cell(0, i)
            cell.text = h
            self._apply_cell_style(cell, is_header=True)

        for r_idx, row in result_df.iterrows():
            if r_idx + 1 >= num_rows:
                print(f"⚠️ 跳过第 {r_idx+1} 行，表格行数不足")
                continue
            for c_idx, value in enumerate(row):
                if c_idx >= num_cols:
                    continue
                try:
                    cell = table.cell(r_idx + 1, c_idx)
                    cell.text = str(value)
                    is_completion_col = headers[c_idx] == "学习完成率"
                    self._apply_cell_style(cell, is_completion_rate=is_completion_col, value=value)
                except Exception as e:
                    print(f"填充单元格 ({r_idx+1}, {c_idx}) 出错: {e}")

        self._add_analysis_textbox(slide, analysis_data, total_info)

        os.makedirs(os.path.dirname(output_path), exist_ok=True)
        prs.save(output_path)
        return output_path

    def _clear_textboxes_containing(self, slide, keywords):
        if not keywords:
            return
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            try:
                text = shape.text_frame.text or ""
            except Exception:
                continue
            if not text:
                continue
            if any(k in text for k in keywords):
                try:
                    shape.text_frame.clear()
                except Exception:
                    try:
                        shape.text_frame.text = ""
                    except Exception:
                        pass

    def _add_analysis_textbox(self, slide, data, total_info):
        left, top, width, height = Inches(7.2), Inches(1.2), Inches(2.6), Inches(5.0)
        tx_box = slide.shapes.add_textbox(left, top, width, height)
        tf = tx_box.text_frame
        tf.word_wrap = True
        try:
            tf.clear()
        except Exception:
            pass

        title_p = tf.paragraphs[0] if tf.paragraphs else tf.add_paragraph()
        title_p.text = "本轮学习情况分析"
        title_p.font.bold = True
        title_p.font.size = Pt(14)
        title_p.font.name = "微软雅黑"

        def to_percent(value):
            if value is None:
                return None
            s = str(value).strip()
            if not s or s == "-":
                return None
            s = s.replace("%", "").replace("％", "").strip()
            try:
                return float(s)
            except Exception:
                return None

        def to_int(value):
            if value is None:
                return None
            s = str(value).strip()
            if not s or s == "-":
                return None
            try:
                return int(float(s))
            except Exception:
                return None

        def fmt_rate(value):
            if value is None:
                return "-"
            if isinstance(value, str) and value.strip().endswith(("%", "％")):
                return value.strip()
            num = to_percent(value)
            if num is None:
                return str(value)
            return f"{num:.2f}%"

        def join_items(items):
            return "、".join([str(x) for x in items if str(x).strip()])

        df = data.copy()
        if "本轮排名" in df.columns:
            df["_rank"] = df["本轮排名"].apply(to_int)
            df = df.sort_values(by=["_rank"], ascending=True, na_position="last").reset_index(drop=True)
        else:
            df["_rate"] = df["学习完成率"].apply(to_percent)
            df = df.sort_values(by=["_rate"], ascending=False, na_position="last").reset_index(drop=True)

        total_count = total_info["总人数"] if isinstance(total_info, dict) else total_info.get("总人数", "-")
        total_rate = total_info["学习完成率"] if isinstance(total_info, dict) else total_info.get("学习完成率", "-")

        lines = []
        lines.append(f"1. 整体情况：本轮总参与人数 {total_count} 人，整体完成率 {total_rate}。")

        if not df.empty:
            top_n = min(3, len(df))
            top_items = [f"{df.iloc[i]['小组']}（{df.iloc[i]['学习完成率']}）" for i in range(top_n)]
            lines.append(f"2. 领先小组：前{top_n}名为 {join_items(top_items)}。")

        if len(df) >= 4:
            bottom_n = min(3, len(df))
            bottom_slice = df.tail(bottom_n)
            bottom_items = [f"{row['小组']}（{row['学习完成率']}）" for _, row in bottom_slice.iterrows()]
            lines.append(f"3. 待提升小组：后{bottom_n}名为 {join_items(bottom_items)}，建议重点跟进未完成人员。")
        elif len(df) >= 2:
            last_row = df.iloc[-1]
            lines.append(f"3. 待提升小组：{last_row['小组']} 当前排名靠后（{last_row['学习完成率']}），建议重点跟进未完成人员。")

        has_prev_rank = "上一轮排名" in df.columns and "本轮排名" in df.columns and any(
            to_int(v) is not None for v in df["上一轮排名"].tolist()
        )
        has_prev_rate = "上一轮学习完成率" in df.columns and any(
            to_percent(v) is not None for v in df["上一轮学习完成率"].tolist()
        )

        next_index = 4
        if has_prev_rank:
            tmp = df.copy()
            tmp["_prev_rank"] = tmp["上一轮排名"].apply(to_int)
            tmp["_curr_rank"] = tmp["本轮排名"].apply(to_int)
            tmp = tmp[tmp["_prev_rank"].notna() & tmp["_curr_rank"].notna()].copy()
            if not tmp.empty:
                tmp["_rank_delta"] = tmp["_prev_rank"] - tmp["_curr_rank"]
                improved = tmp[tmp["_rank_delta"] > 0].sort_values(by="_rank_delta", ascending=False).head(2)
                declined = tmp[tmp["_rank_delta"] < 0].sort_values(by="_rank_delta", ascending=True).head(2)
                if not improved.empty:
                    parts = [f"{r['小组']}（{int(r['_prev_rank'])}→{int(r['_curr_rank'])}）" for _, r in improved.iterrows()]
                    lines.append(f"{next_index}. 排名提升：{join_items(parts)}。")
                    next_index += 1
                if not declined.empty:
                    parts = [f"{r['小组']}（{int(r['_prev_rank'])}→{int(r['_curr_rank'])}）" for _, r in declined.iterrows()]
                    lines.append(f"{next_index}. 排名下滑：{join_items(parts)}。")
                    next_index += 1

        if has_prev_rate:
            tmp = df.copy()
            tmp["_prev_rate"] = tmp["上一轮学习完成率"].apply(to_percent)
            tmp["_curr_rate"] = tmp["学习完成率"].apply(to_percent)
            tmp = tmp[tmp["_prev_rate"].notna() & tmp["_curr_rate"].notna()].copy()
            if not tmp.empty:
                tmp["_rate_delta"] = tmp["_curr_rate"] - tmp["_prev_rate"]
                improved_rate = tmp[tmp["_rate_delta"] > 0].sort_values(by="_rate_delta", ascending=False).head(1)
                declined_rate = tmp[tmp["_rate_delta"] < 0].sort_values(by="_rate_delta", ascending=True).head(1)
                if not improved_rate.empty:
                    r = improved_rate.iloc[0]
                    lines.append(
                        f"{next_index}. 完成率提升：{r['小组']} 提升 {r['_rate_delta']:.2f} 个百分点（{fmt_rate(r['_prev_rate'])}→{fmt_rate(r['_curr_rate'])}）。"
                    )
                    next_index += 1
                if not declined_rate.empty:
                    r = declined_rate.iloc[0]
                    lines.append(
                        f"{next_index}. 完成率下降：{r['小组']} 下降 {abs(r['_rate_delta']):.2f} 个百分点（{fmt_rate(r['_prev_rate'])}→{fmt_rate(r['_curr_rate'])}）。"
                    )
                    next_index += 1

        if not has_prev_rank and not has_prev_rate:
            lines.append("4. 环比变化：上一轮数据不足，本轮结果先作为基线，后续将持续追踪变化。")

        for text in lines:
            p = tf.add_paragraph()
            p.text = text
            p.font.size = Pt(10)
            p.font.name = "微软雅黑"

    def _apply_cell_style(self, cell, is_header=False, is_completion_rate=False, value=None):
        for paragraph in cell.text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            for run in paragraph.runs:
                run.font.size = Pt(10)
                run.font.name = "微软雅黑"
                if is_header:
                    run.font.bold = True

        if is_completion_rate and value is not None and value != "-":
            try:
                rate = float(str(value).replace("%", "").replace("％", ""))
                if rate == 100:
                    color = RGBColor(146, 208, 80)
                elif rate >= 90:
                    color = RGBColor(198, 239, 206)
                elif rate >= 80:
                    color = RGBColor(255, 255, 0)
                else:
                    color = RGBColor(255, 192, 0)
                cell.fill.solid()
                cell.fill.fore_color.rgb = color
            except Exception:
                pass

